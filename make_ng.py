"""
make_ng.py
----------
Создание презентации с фотографиями НГ.

Пользователь присылает Excel-выгрузку (имя не важно). Обрабатывается ПЕРВЫЙ лист.
Фото скачиваются по прямым ссылкам из столбца "Ссылки на фотографии сообщения",
карточки группируются по районам, на выходе — презентация .pptx.

Точка входа для бота:
    process_ng_file(excel_path) -> путь к готовому .pptx

Все промежуточные файлы лежат во временном каталоге запуска (рядом с .pptx),
который вызывающая сторона удаляет после отправки.

Зависимости:
    pip install pandas openpyxl requests python-pptx pillow
"""

import os
import time
import shutil
import tempfile
from io import BytesIO
from pathlib import Path
from datetime import datetime

import pandas as pd
import requests
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Cm, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


# ============================================================
#  ДИЗАЙН
# ============================================================
SLIDE_W = Cm(27.997)
SLIDE_H = Cm(20.999)

HEADER_W = Cm(28.01)
HEADER_H = Cm(1.71)

MARGIN_X = Cm(0.6)
GRID_TOP = HEADER_H + Cm(0.4)
GRID_BOTTOM_MARGIN = Cm(0.8)

COLS, ROWS = 3, 2
MAX_PER_SLIDE = COLS * ROWS
GAP = Cm(0.4)

GRID_W = SLIDE_W - 2 * MARGIN_X
GRID_H = SLIDE_H - GRID_TOP - GRID_BOTTOM_MARGIN

CAPTION_MIN = Cm(2.6)
CAPTION_MAX = Cm(3.6)

C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_DARK   = RGBColor(0x33, 0x33, 0x33)
C_GREY   = RGBColor(0x80, 0x80, 0x80)
C_LIGHT  = RGBColor(0xCC, 0xCC, 0xCC)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_HEADER = RGBColor(0xBF, 0xDF, 0xC5)

LAYOUTS = {
    1: [1],
    2: [2],
    3: [3],
    4: [2, 2],
    5: [3, 2],
    6: [3, 3],
}


def cell_layout(n):
    rows = LAYOUTS.get(n, [3, 3])
    n_rows = len(rows)
    max_cols = max(rows)
    row_h = (GRID_H - GAP * (n_rows - 1)) / n_rows

    caption_h = CAPTION_MAX if n_rows == 1 else CAPTION_MIN
    photo_h = row_h - caption_h

    cell_w = (GRID_W - GAP * (max_cols - 1)) / max_cols

    cells = []
    for r, cols_in_row in enumerate(rows):
        top = GRID_TOP + r * (row_h + GAP)
        row_total_w = cell_w * cols_in_row + GAP * (cols_in_row - 1)
        row_left = MARGIN_X + (GRID_W - row_total_w) / 2
        for c in range(cols_in_row):
            left = row_left + c * (cell_w + GAP)
            cells.append((left, top, cell_w, photo_h, caption_h))
    return cells


# ============================================================
#  СКАЧИВАНИЕ ФОТО
# ============================================================

def parse_photo_urls(photo_column):
    """Разделяет несколько ссылок (через ; или ,) и возвращает первую рабочую."""
    if pd.isna(photo_column) or not str(photo_column).strip():
        return None
    urls = str(photo_column).replace(',', ';').split(';')
    urls = [url.strip() for url in urls if url.strip()]
    if not urls:
        return None
    return urls[0]


def create_placeholder_image(save_path, width=400, height=300):
    """Создаёт заглушку 'Фото нет' если изображение отсутствует."""
    try:
        img = Image.new('RGB', (width, height), color=(220, 220, 220))
        draw = ImageDraw.Draw(img)
        draw.rectangle([0, 0, width - 1, height - 1], outline=(150, 150, 150), width=2)
        try:
            font = ImageFont.truetype("arial.ttf", 24)
        except Exception:
            font = ImageFont.load_default()

        camera_x = width // 2 - 20
        camera_y = height // 2 - 30
        draw.rectangle([camera_x, camera_y, camera_x + 40, camera_y + 30], outline=(100, 100, 100), width=2)
        draw.ellipse([camera_x + 12, camera_y + 8, camera_x + 28, camera_y + 22], outline=(100, 100, 100), width=2)

        bbox = draw.textbbox((0, 0), "Фото нет", font=font)
        text_width = bbox[2] - bbox[0]
        x = (width - text_width) // 2
        y = height // 2 + 20
        draw.text((x, y), "Фото нет", fill=(100, 100, 100), font=font)
        img.save(save_path, quality=95)
        return True
    except Exception as e:
        print(f"Ошибка создания заглушки: {e}")
        return False


def download_image_with_retry(url, save_path, max_retries=3, delay=2):
    """Скачивает изображение с повторными попытками."""
    for attempt in range(max_retries):
        try:
            if url is None:
                return create_placeholder_image(save_path)

            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = requests.get(url, timeout=15, headers=headers)
            response.raise_for_status()

            content_type = response.headers.get('content-type', '')
            if 'image' not in content_type:
                return create_placeholder_image(save_path)

            img = Image.open(BytesIO(response.content))
            if img.mode == 'RGBA':
                img = img.convert('RGB')

            img.thumbnail((1200, 900), Image.Resampling.LANCZOS)
            img.save(save_path, quality=95, optimize=False)
            return True

        except Exception:
            if attempt < max_retries - 1:
                time.sleep(delay)
            else:
                return create_placeholder_image(save_path)
    return False


# ============================================================
#  ОТРИСОВКА СЛАЙДОВ
# ============================================================

def add_rect(slide, left, top, width, height,
             fill=C_WHITE, line=C_LIGHT, line_w=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_header(slide, district_name, photos_in_district):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, HEADER_W, HEADER_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_HEADER
    bar.line.fill.background()
    bar.shadow.inherit = False

    tb_w = Cm(14.0)
    tb_h = HEADER_H - Cm(0.2)
    tb_left = SLIDE_W - MARGIN_X - tb_w
    tb_top = Cm(0.1)
    tb = slide.shapes.add_textbox(tb_left, tb_top, tb_w, tb_h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.RIGHT
    r1 = p1.add_run()
    r1.text = f"Район: {district_name}"
    r1.font.name = "Calibri"
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.color.rgb = C_BLACK

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"Фото в районе: {photos_in_district}"
    r2.font.name = "Calibri"
    r2.font.size = Pt(11)
    r2.font.color.rgb = C_DARK


def add_footer(slide, district_name):
    fw = Cm(12.0)
    fh = Cm(0.5)
    tb = slide.shapes.add_textbox(SLIDE_W - MARGIN_X - fw, SLIDE_H - Cm(0.6), fw, fh)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"Район: {district_name}"
    r.font.name = "Calibri"
    r.font.size = Pt(8)
    r.font.color.rgb = C_GREY


def create_photo_card(slide, image_path, data, cell):
    left, top, cell_w, photo_h, caption_h = cell

    if image_path and os.path.exists(image_path):
        try:
            pic = slide.shapes.add_picture(
                image_path, left, top,
                width=int(cell_w), height=int(photo_h),
            )
            pic.line.color.rgb = C_LIGHT
            pic.line.width = Pt(0.5)
        except Exception as e:
            print(f"Ошибка добавления изображения {image_path}: {e}")
            add_rect(slide, left, top, cell_w, photo_h,
                     fill=RGBColor(0xF5, 0xF5, 0xF5), line=C_LIGHT, line_w=0.5)
    else:
        add_rect(slide, left, top, cell_w, photo_h,
                 fill=RGBColor(0xF5, 0xF5, 0xF5), line=C_LIGHT, line_w=0.5)
        tb = slide.shapes.add_textbox(left, top, cell_w, photo_h)
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "Фото нет"
        r.font.name = "Calibri"
        r.font.size = Pt(14)
        r.font.color.rgb = C_GREY

    cap_top = top + photo_h + Cm(0.1)
    tb = slide.shapes.add_textbox(left, cap_top, cell_w, caption_h - Cm(0.1))
    tf = tb.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.word_wrap = True

    ticket  = data.get('Номер сообщения', 'Не указан')
    address = data.get('Адрес', 'Не указан')
    problem = data.get('Проблемная тема', 'Не указана')
    date    = data.get('Дата публикации заявки', 'Не указана')

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r = p1.add_run()
    r.text = f"№ {ticket}   "
    r.font.name = "Calibri"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = C_BLACK
    r = p1.add_run()
    r.text = str(date)
    r.font.name = "Calibri"
    r.font.size = Pt(12)
    r.font.color.rgb = C_DARK

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(2)
    r = p2.add_run()
    r.text = str(address)
    r.font.name = "Calibri"
    r.font.size = Pt(11)
    r.font.color.rgb = C_DARK

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.LEFT
    p3.space_before = Pt(2)
    r = p3.add_run()
    r.text = str(problem)
    r.font.name = "Calibri"
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = C_GREY


def create_district_slide(prs, district_name, district_data_with_paths,
                          slide_num_in_district, photos_in_district):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, district_name, photos_in_district)

    total_cards = len(district_data_with_paths)
    start_idx = slide_num_in_district * MAX_PER_SLIDE
    end_idx = min(start_idx + MAX_PER_SLIDE, total_cards)
    cards_on_slide = end_idx - start_idx

    cells = cell_layout(cards_on_slide)

    for card_pos in range(cards_on_slide):
        data_idx = start_idx + card_pos
        row_data = district_data_with_paths.iloc[data_idx]

        card_data = {
            "Номер сообщения": row_data["Номер сообщения"] if pd.notna(row_data["Номер сообщения"]) else "Не указан",
            "Адрес": str(row_data["Адрес"]) if pd.notna(row_data["Адрес"]) else "Не указан",
            "Проблемная тема": str(row_data["Проблемная тема"]) if pd.notna(row_data["Проблемная тема"]) else "Не указана",
            "Дата публикации заявки": row_data["Дата публикации заявки"] if pd.notna(row_data["Дата публикации заявки"]) else "Не указана",
        }

        img_path = row_data.get('image_path', None)
        create_photo_card(slide, img_path, card_data, cells[card_pos])

    add_footer(slide, district_name)
    return slide


# ============================================================
#  ОСНОВНОЙ ПАЙПЛАЙН
# ============================================================

def process_excel_to_ppt(excel_path, output_ppt_path, images_temp_dir="temp_images"):
    """Читает Excel (первый лист) и создаёт PPT с карточками по районам."""
    print(f"Чтение Excel файла: {excel_path}")
    try:
        df = pd.read_excel(excel_path)       # первый лист
        print(f"Найдено {len(df)} записей")
    except Exception as e:
        raise RuntimeError(f"Ошибка чтения Excel файла: {e}")

    required_columns = ["Номер сообщения", "Адрес", "Проблемная тема",
                        "Дата публикации заявки", "Ссылки на фотографии сообщения", "Район"]

    missing_columns = [col for col in required_columns if col not in df.columns]
    if missing_columns:
        raise ValueError(
            f"Отсутствуют столбцы: {missing_columns}. "
            f"Доступные: {list(df.columns)}"
        )

    df['Район'] = df['Район'].fillna('Не указан')
    districts = df['Район'].unique()
    print(f"Найдено районов: {len(districts)}")

    Path(images_temp_dir).mkdir(parents=True, exist_ok=True)

    success_count = 0
    placeholder_count = 0
    df['image_path'] = None

    for idx, row in df.iterrows():
        photo_url = parse_photo_urls(row["Ссылки на фотографии сообщения"])
        local_path = os.path.join(images_temp_dir, f"image_{idx}.jpg")

        if download_image_with_retry(photo_url, local_path, max_retries=2, delay=1):
            df.at[idx, 'image_path'] = local_path
            if os.path.getsize(local_path) < 1000:
                placeholder_count += 1
            else:
                success_count += 1
        else:
            placeholder_count += 1

    print(f"Статистика: {success_count} фото, {placeholder_count} заглушек")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total_slides = 0
    for district in districts:
        district_data = df[df['Район'] == district].copy().reset_index(drop=True)
        photos_in_district = len(district_data)

        num_slides_for_district = (photos_in_district + MAX_PER_SLIDE - 1) // MAX_PER_SLIDE
        for slide_num in range(num_slides_for_district):
            create_district_slide(prs, district, district_data,
                                  slide_num, photos_in_district)
            total_slides += 1

    prs.save(output_ppt_path)
    print(f"Готово: {output_ppt_path}  (карточек: {len(df)}, "
          f"слайдов: {total_slides}, районов: {len(districts)})")
    return True


# ============================================================
#  ТОЧКА ВХОДА ДЛЯ БОТА
# ============================================================

def process_ng_file(excel_path: str) -> str:
    """
    Полный пайплайн НГ. Возвращает путь к готовой презентации .pptx.
    Все промежуточные файлы лежат в одном временном каталоге (рядом с .pptx),
    который вызывающая сторона удаляет после отправки.
    При ошибке временный каталог очищается, исключение пробрасывается.
    """
    run_dir = Path(tempfile.mkdtemp(prefix="ng_photo_"))
    try:
        images_dir = str(run_dir / "temp_images")
        ts = datetime.now().strftime("%d.%m.%Y_%H-%M")
        output_pptx = str(run_dir / f"Фото_НГ_{ts}.pptx")

        process_excel_to_ppt(excel_path, output_pptx, images_dir)
        return output_pptx
    except Exception:
        shutil.rmtree(run_dir, ignore_errors=True)
        raise


if __name__ == "__main__":
    import sys
    test_path = sys.argv[1] if len(sys.argv) > 1 else "1.xlsx"
    out = process_ng_file(test_path)
    print("Готово:", out)
