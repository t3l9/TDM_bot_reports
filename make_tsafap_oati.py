"""
make_tsafap_oati.py
-------------------
Создание презентации с фотографиями ОАТИ / ЦАФАП.

Один и тот же модуль обслуживает два режима (system="ОАТИ" или "ЦАФАП"):
отличаются только подпись в шапке и заголовок сводной таблицы.

Пайплайн:
  1) Selenium открывает каждую ссылку из столбца "Скриншоты" выгрузки
     и скачивает фото-скриншоты во временную папку photos/{ID предписания}/
     (с автоматическим сжатием, чтобы .pptx не раздувался).
  2) Слайд со сводной таблицей по районам (Балансодержатель × Тип объекта).
  3) Слайды с фотографиями, сгруппированные по районам.

Точка входа для бота:
    process_tsafap_oati_file(excel_path, system="ЦАФАП") -> путь к .pptx

Особенности:
  * Имя входного файла НЕ важно.
  * Лист берётся "Детализированный", если он есть, иначе — первый лист книги.
  * Герб ЮВАО берётся из файла gerb_uvao.png рядом с модулем (если нет — пропускается).

Зависимости:
    pip install selenium pandas openpyxl requests python-pptx pillow
"""

import os
import re
import time
import shutil
import zipfile
import tempfile
import pandas as pd
import requests
from io import BytesIO
from pathlib import Path
from datetime import datetime, date, timedelta
from urllib.parse import urljoin

from PIL import Image

# Selenium
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.common.exceptions import TimeoutException, NoSuchElementException

# PPTX
from pptx import Presentation
from pptx.util import Inches, Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ============================================================
#  ОБЩИЕ НАСТРОЙКИ
# ============================================================
SHEET_NAME_PREFERRED = "Детализированный"   # если листа нет — берём первый

# Selenium
HEADLESS        = True
PAGE_LOAD_WAIT  = 15
CAROUSEL_WAIT   = 2.5
MAX_PHOTOS      = 30
SKIP_EXISTING   = True
CHROMEDRIVER_PATH = None
BASE_URL = "https://cafap.mos.ru"

# Сжатие скачанных скриншотов (иначе огромные PNG раздувают .pptx)
MAX_IMG_W    = 1280
MAX_IMG_H    = 960
JPEG_QUALITY = 80

# Файл герба префектуры ЮВАО (PNG, желательно с прозрачным фоном).
# Лежит рядом с этим модулем. Если файла нет — логотип просто не вставляется.
GERB_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "gerb_uvao.png")
GERB_HEIGHT_CM = 1.15   # высота герба в шапке (см)

# Столбцы выгрузки, используемые в сводной таблице
DISTRICT_COLUMN = "Район"
TYPE_COLUMN     = "Тип объекта (АСУ ОДС)"

# Фиксированный порядок балансодержателей (районов) ЮВАО для таблицы.
# Районы из выгрузки, которых здесь нет, добавляются в конец.
DISTRICTS_UVAO = [
    "Выхино-Жулебино", "Капотня", "Кузьминки", "Лефортово", "Люблино",
    "Марьино", "Некрасовка", "Нижегородский", "Печатники", "Рязанский",
    "Текстильщики", "Южнопортовый", "АВД ЮВАО",
]

# (Подсветка «Итого» считается автоматически: топ-3 значений с учётом ничьих.)

# Конфигурация режимов
SYSTEM_CONFIG = {
    "ЦАФАП": {
        "brand": "Система «ЦАФАП»",
        "short": "ЦАФАП",
        "table_title": "Сводная информация по нарушениям, выявленных информационной системой «ЦАФАП»",
    },
    "ОАТИ": {
        "brand": "Система «ОАТИ»",
        "short": "ОАТИ",
        "table_title": "Сводная информация по нарушениям, выявленных «ОАТИ»",
    },
}


# ============================================================
#  ЦВЕТА / ГЕОМЕТРИЯ
# ============================================================
C_BLACK    = RGBColor(0x00, 0x00, 0x00)
C_DARK     = RGBColor(0x33, 0x33, 0x33)
C_GREY     = RGBColor(0x80, 0x80, 0x80)
C_LIGHT    = RGBColor(0xCC, 0xCC, 0xCC)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_HEADER   = RGBColor(0xBF, 0xDF, 0xC5)
C_ORANGE   = RGBColor(0xE2, 0x6A, 0x2C)   # «УИС»
C_NAVY     = RGBColor(0x1F, 0x38, 0x64)   # текст в таблице
C_TBL_HEAD = RGBColor(0xB4, 0xC6, 0xE7)   # шапка таблицы (голубой)
C_TBL_RED  = RGBColor(0xFF, 0x00, 0x00)   # подсветка итога

SLIDE_W = Cm(27.997)
SLIDE_H = Cm(20.999)

HEADER_W = Cm(28.01)
HEADER_H = Cm(1.71)

MARGIN_X = Cm(0.6)
GRID_TOP = HEADER_H + Cm(0.4)
GRID_BOTTOM_MARGIN = Cm(0.8)

GAP = Cm(0.4)

GRID_W = SLIDE_W - 2 * MARGIN_X
GRID_H = SLIDE_H - GRID_TOP - GRID_BOTTOM_MARGIN

MAX_PER_SLIDE = 6
CAPTION_MIN = Cm(2.6)


def _emu(v):
    """Округляет координату/размер до ЦЕЛОГО EMU.
    Дробные EMU (например 720000.0 после деления) недопустимы по схеме OOXML —
    PowerPoint считает такой файл повреждённым и «чинит» его."""
    return Emu(int(round(float(v))))

# Точные размеры ФОТО (ширина, высота) в см для слайдов с < 6 фото
EXPLICIT_SIZES = {
    1: (17.78, 9.61),
    2: (13.20, 9.38),
    3: (8.67,  9.16),
    4: (10.09, 6.25),
}
EXPLICIT_CAPTION = {   # высота подписи под фото, см
    1: 3.2,
    2: 3.2,
    3: 3.0,
    4: 2.6,
}


# ============================================================
#  РАСКЛАДКА ФОТО НА СЛАЙДЕ
# ============================================================

def _grid_layout(n):
    """Сетка 3 колонки (для 5 и 6 фото). 5 → 3 сверху + 2 по центру снизу."""
    rows = {5: [3, 2], 6: [3, 3]}.get(n, [3, 3])
    n_rows = len(rows)
    row_h = (GRID_H - GAP * (n_rows - 1)) / n_rows
    caption_h = CAPTION_MIN
    photo_h = row_h - caption_h
    cell_w = (GRID_W - GAP * 2) / 3

    cells = []
    for r, cols in enumerate(rows):
        top = GRID_TOP + r * (row_h + GAP)
        row_w = cell_w * cols + GAP * (cols - 1)
        left0 = MARGIN_X + (GRID_W - row_w) / 2   # центрируем ряд
        for c in range(cols):
            left = left0 + c * (cell_w + GAP)
            cells.append((left, top, cell_w, photo_h, caption_h))
    return cells


def cell_layout(n):
    """Возвращает список ячеек (left, top, photo_w, photo_h, caption_h)."""
    if n <= 0:
        return []

    if n in EXPLICIT_SIZES:
        pw = Cm(EXPLICIT_SIZES[n][0])
        ph = Cm(EXPLICIT_SIZES[n][1])
        cap_h = Cm(EXPLICIT_CAPTION[n])
        rows = [n] if n <= 3 else [2, 2]   # 4 → 2×2
        n_rows = len(rows)

        block_h = ph + Cm(0.1) + cap_h
        total_h = block_h * n_rows + GAP * (n_rows - 1)
        top0 = GRID_TOP + (GRID_H - total_h) / 2
        if top0 < GRID_TOP:
            top0 = GRID_TOP

        cells = []
        for r, cols in enumerate(rows):
            top = top0 + r * (block_h + GAP)
            row_w = pw * cols + GAP * (cols - 1)
            left0 = (SLIDE_W - row_w) / 2        # центрируем по слайду
            for c in range(cols):
                left = left0 + c * (pw + GAP)
                cells.append((left, top, pw, ph, cap_h))
        return cells

    # 5, 6 (и любой fallback) — сетка
    return _grid_layout(min(n, 6))


# ============================================================
#  ФАЗА 1: СКАЧИВАНИЕ ФОТО (Selenium)
# ============================================================

def make_driver(headless: bool = HEADLESS):
    opts = Options()
    if headless:
        opts.add_argument("--headless=new")
    opts.add_argument("--window-size=1400,900")
    opts.add_argument("--disable-blink-features=AutomationControlled")
    opts.add_experimental_option("excludeSwitches", ["enable-automation"])
    opts.add_experimental_option("useAutomationExtension", False)
    if CHROMEDRIVER_PATH:
        service = Service(executable_path=CHROMEDRIVER_PATH)
        return webdriver.Chrome(service=service, options=opts)
    return webdriver.Chrome(options=opts)


def get_current_img_uuid(driver):
    try:
        img = driver.find_element(By.CSS_SELECTOR, "div.img-wrapper img")
        src = img.get_attribute("src") or ""
    except NoSuchElementException:
        return None
    m = re.search(r"/api/issue/screenshot/get/([0-9a-fA-F\-]+)", src)
    return m.group(1) if m else None


def click_next(driver):
    selectors = [
        "button.carousel-next", ".carousel-next",
        "button[aria-label*='next' i]", "button[aria-label*='след' i]",
        ".swiper-button-next", ".carousel-control-next",
        "[class*='next']:not([disabled])",
    ]
    for sel in selectors:
        try:
            els = driver.find_elements(By.CSS_SELECTOR, sel)
            for el in els:
                if el.is_displayed() and el.is_enabled():
                    driver.execute_script("arguments[0].click();", el)
                    return True
        except Exception:
            continue
    try:
        body = driver.find_element(By.TAG_NAME, "body")
        body.send_keys(Keys.ARROW_RIGHT)
        return True
    except Exception:
        return False


def collect_uuids(driver, page_url):
    driver.get(page_url)
    wait = WebDriverWait(driver, PAGE_LOAD_WAIT)
    try:
        wait.until(EC.presence_of_element_located(
            (By.CSS_SELECTOR, "div.img-wrapper img")
        ))
    except TimeoutException:
        return []

    seen = []
    uuid = get_current_img_uuid(driver)
    if uuid:
        seen.append(uuid)

    for _ in range(MAX_PHOTOS - 1):
        prev = uuid
        if not click_next(driver):
            break
        time.sleep(CAROUSEL_WAIT)
        uuid = get_current_img_uuid(driver)
        if not uuid:
            break
        if uuid in seen or uuid == prev:
            break
        seen.append(uuid)
    return seen


def cookies_for_requests(driver):
    return {c["name"]: c["value"] for c in driver.get_cookies()}


def download_image(uuid, cookies, dest_path, base_url):
    """Скачивает скриншот и сохраняет ТОЛЬКО как валидный JPEG.
    Если содержимое не открывается как картинка (HTML-ошибка, webp без поддержки
    и т.п.) — возвращает None и ничего не пишет, чтобы в .pptx не попал битый файл,
    из-за которого PowerPoint «чинит» презентацию."""
    url = urljoin(base_url, f"/api/issue/screenshot/get/{uuid}")
    headers = {"User-Agent": "Mozilla/5.0", "Referer": base_url}
    r = requests.get(url, cookies=cookies, headers=headers, timeout=30)
    r.raise_for_status()

    try:
        img = Image.open(BytesIO(r.content))
        img.load()
        if img.mode != "RGB":
            img = img.convert("RGB")
        img.thumbnail((MAX_IMG_W, MAX_IMG_H), Image.Resampling.LANCZOS)
        final = dest_path.with_suffix(".jpg")
        img.save(final, format="JPEG", quality=JPEG_QUALITY, optimize=True)
        # контрольная проверка, что файл читается обратно
        with Image.open(final) as chk:
            chk.verify()
        return final
    except Exception as e:
        print(f"   ! пропускаю фото (не удалось обработать): {e}")
        return None


def download_phase(df, photos_dir: Path, headless: bool = HEADLESS):
    df = df[df["Скриншоты"].notna()].copy()
    df["ID предписания"] = df["ID предписания"].astype(str)
    total = len(df)
    print(f"=== ФАЗА 1: скачивание фото === заявок: {total}")

    photos_dir.mkdir(parents=True, exist_ok=True)
    driver = make_driver(headless)

    ok = 0
    try:
        for i, (_, row) in enumerate(df.iterrows(), 1):
            issue_id = str(row["ID предписания"]).strip()
            url = str(row["Скриншоты"]).strip()
            out_dir = photos_dir / issue_id

            if SKIP_EXISTING and out_dir.exists() and any(out_dir.iterdir()):
                ok += 1
                continue

            out_dir.mkdir(parents=True, exist_ok=True)
            print(f"[{i}/{total}] {issue_id}: {url}")

            try:
                uuids = collect_uuids(driver, url)
            except Exception as e:
                print(f"   ! ошибка сбора UUID: {e}")
                continue

            if not uuids:
                print("   ! фото не найдены на странице")
                continue

            cookies = cookies_for_requests(driver)
            saved = 0
            for n, uuid in enumerate(uuids, 1):
                dst = out_dir / f"{n:02d}"
                try:
                    res = download_image(uuid, cookies, dst, BASE_URL)
                    if res:
                        saved += 1
                except Exception as e:
                    print(f"   ! не скачалось {uuid[:8]}…: {e}")
            if saved:
                ok += 1
    finally:
        driver.quit()
    print(f"Скачано/найдено фото для {ok} из {total} заявок.")
    return ok


# ============================================================
#  ШАПКА / БРЕНДИНГ (на всех слайдах)
# ============================================================

def _add_header_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, _emu(HEADER_W), _emu(HEADER_H))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_HEADER
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(_emu(left), _emu(top), _emu(width), _emu(height))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    return tb, tf


def _run(p, text, size, bold=False, color=C_BLACK, italic=False, name="Calibri"):
    r = p.add_run()
    r.text = text
    r.font.name = name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def add_brand(slide, brand_label: str):
    """Левая часть шапки: герб + «Префектура ЮВАО» | «УИС» + Система «...»."""
    # Герб
    logo_right = Cm(0.3)
    if GERB_PATH and os.path.exists(GERB_PATH):
        try:
            gh = Cm(GERB_HEIGHT_CM)
            gtop = _emu((HEADER_H - gh) / 2)
            pic = slide.shapes.add_picture(GERB_PATH, Cm(0.3), gtop, height=gh)
            logo_right = Cm(0.3) + pic.width
        except Exception as e:
            print(f"   ! не удалось вставить герб: {e}")
            logo_right = Cm(1.5)
    else:
        logo_right = Cm(0.3)

    x = logo_right + Cm(0.15)

    # «Префектура ЮВАО» (две строки)
    tb, tf = _textbox(slide, x, Cm(0.18), Cm(2.3), Cm(1.35))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    _run(p1, "Префектура", 9, bold=True, color=C_DARK)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    _run(p2, "ЮВАО", 9, bold=True, color=C_DARK)

    x = x + Cm(2.45)

    # Разделитель
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(x), Cm(0.28), Cm(0.03), Cm(1.15))
    div.fill.solid()
    div.fill.fore_color.rgb = C_GREY
    div.line.fill.background()
    div.shadow.inherit = False

    x = x + Cm(0.2)

    # «УИС»
    tb, tf = _textbox(slide, x, Cm(0.3), Cm(1.2), Cm(1.1))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _run(p, "УИС", 12, bold=True, color=C_ORANGE)

    x = x + Cm(1.35)

    # Система «...»
    tb, tf = _textbox(slide, x, Cm(0.1), Cm(11.0), Cm(1.5))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _run(p, brand_label, 18, bold=True, color=C_BLACK)


def add_header(slide, brand_label, district, photos_in_district):
    """Полная шапка фото-слайда: бренд слева + район/кол-во справа."""
    _add_header_bar(slide)
    add_brand(slide, brand_label)

    tb_w = Cm(10.0)
    tb_left = SLIDE_W - MARGIN_X - tb_w
    tb, tf = _textbox(slide, tb_left, Cm(0.1), tb_w, HEADER_H - Cm(0.2))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.RIGHT
    _run(p1, f"Район: {district}", 16, bold=True, color=C_BLACK)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    _run(p2, f"Фото в районе: {photos_in_district}", 11, color=C_DARK)


# ============================================================
#  ФОТО-СЛАЙДЫ
# ============================================================

def fmt_date(v):
    if pd.isna(v):
        return ""
    if isinstance(v, (datetime, date)):
        return v.strftime("%d.%m.%Y")
    s = str(v).strip()
    s = re.sub(r"\s+\d{1,2}:\d{2}(:\d{2})?$", "", s)
    return s


def safe(v):
    return "" if pd.isna(v) else str(v).strip()


def list_photos(issue_id, photos_dir: Path):
    d = photos_dir / str(issue_id).strip()
    if not d.exists():
        return []
    exts = {".jpg", ".jpeg", ".png", ".webp", ".bmp"}
    return sorted(p for p in d.iterdir() if p.suffix.lower() in exts)


def add_rect(slide, left, top, width, height, fill=C_WHITE, line=C_LIGHT, line_w=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(left), _emu(top),
                                   _emu(width), _emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=10, bold=False, color=C_BLACK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(_emu(left), _emu(top), _emu(width), _emu(height))
    tf = tb.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    _run(p, text, font_size, bold=bold, color=color)
    return tb


def _valid_image(path) -> bool:
    """True, если файл — целое изображение (чтобы не вставлять битое в .pptx)."""
    try:
        with Image.open(str(path)) as im:
            im.verify()
        return True
    except Exception:
        return False


def place_photo(slide, cell, photo_path, item):
    left, top, photo_w, photo_h, caption_h = cell

    if photo_path and photo_path.exists() and _valid_image(photo_path):
        try:
            pic = slide.shapes.add_picture(str(photo_path), _emu(left), _emu(top),
                                           width=_emu(photo_w), height=_emu(photo_h))
            pic.line.color.rgb = C_LIGHT
            pic.line.width = Pt(0.5)
        except Exception as e:
            add_rect(slide, left, top, photo_w, photo_h,
                     fill=RGBColor(0xF5, 0xF5, 0xF5), line=C_LIGHT, line_w=0.5)
            add_textbox(slide, left, top, photo_w, photo_h, "[ошибка фото]",
                        font_size=9, color=C_GREY, align=PP_ALIGN.CENTER,
                        anchor=MSO_ANCHOR.MIDDLE)
    else:
        add_rect(slide, left, top, photo_w, photo_h,
                 fill=RGBColor(0xF5, 0xF5, 0xF5), line=C_LIGHT, line_w=0.5)
        add_textbox(slide, left, top, photo_w, photo_h, "[фото отсутствует]",
                    font_size=10, color=C_GREY, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)

    cap_top = top + photo_h + Cm(0.1)
    tb = slide.shapes.add_textbox(_emu(left), _emu(cap_top), _emu(photo_w),
                                  _emu(caption_h - Cm(0.1)))
    tf = tb.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _run(p, f"№ {item['id']}   ", 12, bold=True, color=C_BLACK)
    _run(p, item["date"], 12, color=C_DARK)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(2)
    _run(p2, item["address"], 11, color=C_DARK)

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.LEFT
    p3.space_before = Pt(2)
    _run(p3, item["violation"], 10, italic=True, color=C_GREY)


def add_title_slide(prs, brand_label, total_districts, total_items, total_photos):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide)
    add_brand(slide, brand_label)

    add_textbox(slide, Cm(1.0), Cm(7.0), SLIDE_W - Cm(2.0), Cm(2.0),
                "Отчёт по нарушениям", font_size=36, bold=True, color=C_BLACK)
    add_textbox(slide, Cm(1.0), Cm(9.5), SLIDE_W - Cm(2.0), Cm(1.0),
                "Группировка по районам", font_size=18, color=C_DARK)
    add_textbox(slide, Cm(1.0), Cm(12.0), SLIDE_W - Cm(2.0), Cm(1.0),
                f"Районов: {total_districts}    Нарушений: {total_items}    Фото: {total_photos}",
                font_size=14, color=C_GREY)
    add_textbox(slide, Cm(1.0), SLIDE_H - Cm(1.2), SLIDE_W - Cm(2.0), Cm(0.8),
                datetime.now().strftime("Сформировано %d.%m.%Y"),
                font_size=10, color=C_GREY)


def add_district_divider(prs, brand_label, district, items_count, photos_count):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide)
    add_brand(slide, brand_label)

    band_h = Cm(4.0)
    band_top = (SLIDE_H - band_h) / 2
    add_rect(slide, 0, band_top, SLIDE_W, band_h, fill=C_BLACK, line=C_BLACK)
    add_textbox(slide, MARGIN_X, band_top + Cm(0.4), SLIDE_W - 2 * MARGIN_X, Cm(2.2),
                district, font_size=36, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, MARGIN_X, band_top + Cm(2.6), SLIDE_W - 2 * MARGIN_X, Cm(1.0),
                f"Нарушений: {items_count}     Фото: {photos_count}",
                font_size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)


# ============================================================
#  СВОДНАЯ ТАБЛИЦА
# ============================================================

def _yesterday_str():
    return (datetime.now() - timedelta(days=1)).strftime("%d.%m.%Y")


def _compute_summary(df):
    """Возвращает (types, rows, totals_per_type, grand_total).
    rows: список (district, [count_по_типам], total)."""
    if DISTRICT_COLUMN not in df.columns:
        return None

    sub = df.copy()
    sub[DISTRICT_COLUMN] = sub[DISTRICT_COLUMN].fillna("Не указан").astype(str).str.strip()

    if TYPE_COLUMN in df.columns:
        sub[TYPE_COLUMN] = sub[TYPE_COLUMN].fillna("Не указан").astype(str).str.strip()
        types = sorted(sub[TYPE_COLUMN].unique(), key=lambda s: s.lower())
        counts = sub.groupby([DISTRICT_COLUMN, TYPE_COLUMN]).size().to_dict()
        def cell(d, t): return int(counts.get((d, t), 0))
    else:
        types = ["Всего"]
        counts = sub.groupby([DISTRICT_COLUMN]).size().to_dict()
        def cell(d, t): return int(counts.get(d, 0))

    present = list(dict.fromkeys(sub[DISTRICT_COLUMN].tolist()))
    districts = list(DISTRICTS_UVAO) + [d for d in present if d not in DISTRICTS_UVAO]

    rows = []
    totals_per_type = [0] * len(types)
    grand_total = 0
    for d in districts:
        vals = [cell(d, t) for t in types]
        total = sum(vals)
        rows.append((d, vals, total))
        for i, v in enumerate(vals):
            totals_per_type[i] += v
        grand_total += total

    return types, rows, totals_per_type, grand_total


def _cell_set(cell, text, size=11, bold=False, color=C_NAVY,
              align=PP_ALIGN.CENTER, fill=None):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_WHITE
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Cm(0.1)
    cell.margin_right = Cm(0.1)
    cell.margin_top = Cm(0.02)
    cell.margin_bottom = Cm(0.02)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _run(p, str(text), size, bold=bold, color=color)


def _cell_borders(cell, color="404040", width_pt=0.75):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    w = str(int(width_pt * 12700))
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for e in tcPr.findall(qn(tag)):
            tcPr.remove(e)
    for tag in ("a:lnB", "a:lnT", "a:lnR", "a:lnL"):   # вставляем в начало → порядок L,R,T,B
        ln = etree.Element(qn(tag))
        ln.set("w", w); ln.set("cap", "flat"); ln.set("cmpd", "sng"); ln.set("algn", "ctr")
        fill = etree.SubElement(ln, qn("a:solidFill"))
        clr = etree.SubElement(fill, qn("a:srgbClr"))
        clr.set("val", color)
        tcPr.insert(0, ln)


def build_summary_table_slide(prs, df, cfg):
    res = _compute_summary(df)
    if not res:
        print("   ! нет столбца 'Район' — сводная таблица пропущена")
        return
    types, rows, totals_per_type, grand_total = res
    short = cfg["short"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide)
    add_brand(slide, cfg["brand"])

    # Заголовок над таблицей
    tb, tf = _textbox(slide, Cm(1.0), Cm(2.3), SLIDE_W - Cm(2.0), Cm(1.6))
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    _run(p1, cfg["table_title"], 13, bold=True, color=C_BLACK)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    _run(p2, f"По состоянию за {_yesterday_str()}", 13, bold=True, color=C_BLACK)

    # Геометрия таблицы
    nt = len(types)
    num_w, bal_w, tot_w = 1.0, 5.2, 2.6
    avail = 24.0 - (num_w + bal_w + tot_w)
    type_w = max(2.2, min(3.6, avail / max(nt, 1)))
    table_w_cm = num_w + bal_w + tot_w + type_w * nt

    n_rows = 2 + len(rows) + 1
    n_cols = nt + 3
    left = (SLIDE_W - Cm(table_w_cm)) / 2

    # Высота таблицы для строгого центрирования по слайду
    h0, h1, hd, hr = 0.95, 0.85, 0.72, 0.8
    table_h_cm = h0 + h1 + hd * len(rows) + hr
    top = max(Cm(4.2), (SLIDE_H - Cm(table_h_cm)) / 2)

    gfx = slide.shapes.add_table(n_rows, n_cols, _emu(left), _emu(top),
                                 _emu(Cm(table_w_cm)), _emu(Cm(table_h_cm)))
    table = gfx.table
    table.first_row = False
    table.horz_banding = False

    # Ширины колонок
    table.columns[0].width = _emu(Cm(num_w))
    table.columns[1].width = _emu(Cm(bal_w))
    for j in range(nt):
        table.columns[2 + j].width = _emu(Cm(type_w))
    table.columns[n_cols - 1].width = _emu(Cm(tot_w))

    # Высоты строк
    table.rows[0].height = _emu(Cm(h0))
    table.rows[1].height = _emu(Cm(h1))
    for r in range(2, n_rows - 1):
        table.rows[r].height = _emu(Cm(hd))
    table.rows[n_rows - 1].height = _emu(Cm(hr))

    last = n_cols - 1

    # ── Шапка таблицы ──
    table.cell(0, 0).merge(table.cell(1, 0))
    table.cell(0, 1).merge(table.cell(1, 1))
    if nt > 1:
        table.cell(0, 2).merge(table.cell(0, 2 + nt - 1))
    table.cell(0, last).merge(table.cell(1, last))

    _cell_set(table.cell(0, 0), "№", 12, bold=True, fill=C_TBL_HEAD)
    _cell_set(table.cell(0, 1), "Балансодержатель", 12, bold=True, fill=C_TBL_HEAD)
    _cell_set(table.cell(0, 2), f"Всего нарушений {short}", 12, bold=True, fill=C_TBL_HEAD)
    _cell_set(table.cell(0, last), f"Итого {short}", 11, bold=True, fill=C_TBL_HEAD)
    for j, t in enumerate(types):
        _cell_set(table.cell(1, 2 + j), t, 11, bold=True, fill=C_TBL_HEAD)

    # ── Данные ──
    # Топ-3 значений «Итого» (по различным величинам, с учётом ничьих).
    # Например для 33,22,22,17,17 различные = {33,22,17} → красим все эти строки.
    distinct_desc = sorted({t for (_, _, t) in rows if t > 0}, reverse=True)
    red_threshold = distinct_desc[2] if len(distinct_desc) >= 3 else (
        distinct_desc[-1] if distinct_desc else None)

    def _is_red(total):
        return total > 0 and red_threshold is not None and total >= red_threshold

    for i, (district, vals, total) in enumerate(rows):
        r = 2 + i
        _cell_set(table.cell(r, 0), i + 1, 11, color=C_NAVY)
        _cell_set(table.cell(r, 1), district, 11, color=C_NAVY, align=PP_ALIGN.LEFT)
        for j, v in enumerate(vals):
            _cell_set(table.cell(r, 2 + j), v, 11, color=C_NAVY)
        red = _is_red(total)
        _cell_set(table.cell(r, last), total, 11, bold=red,
                  color=(C_WHITE if red else C_NAVY),
                  fill=(C_TBL_RED if red else None))

    # ── Итого по округу ──
    rt = n_rows - 1
    table.cell(rt, 0).merge(table.cell(rt, 1))
    _cell_set(table.cell(rt, 0), "ИТОГО ПО ОКРУГУ:", 11, bold=True,
              color=C_NAVY, fill=C_TBL_HEAD, align=PP_ALIGN.CENTER)
    for j, v in enumerate(totals_per_type):
        _cell_set(table.cell(rt, 2 + j), v, 11, bold=True, color=C_NAVY, fill=C_TBL_HEAD)
    _cell_set(table.cell(rt, last), grand_total, 11, bold=True, color=C_NAVY, fill=C_TBL_HEAD)

    # Границы для всех ячеек, КРОМЕ поглощённых ячеек объединения
    # (применение tcPr к hMerge/vMerge-ячейкам ломает файл в PowerPoint).
    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            if not cell.is_spanned:
                _cell_borders(cell)


# ============================================================
#  СБОРКА ПРЕЗЕНТАЦИИ
# ============================================================

def build_phase(df, photos_dir: Path, output_pptx: str, cfg: dict):
    print("=== ФАЗА 2: сборка презентации ===")

    cards_by_district = {}
    total_items = 0
    total_photos = 0

    df_sorted = df.sort_values(
        by=["Район", "Дата фиксации нарушения"],
        kind="stable", na_position="last",
    ) if "Дата фиксации нарушения" in df.columns else df

    for _, row in df_sorted.iterrows():
        district = safe(row.get("Район")) or "Без района"
        item = {
            "id":        safe(row.get("ID предписания")),
            "date":      fmt_date(row.get("Дата фиксации нарушения")),
            "address":   safe(row.get("Наименование объекта")) or "—",
            "violation": safe(row.get("Наименование нарушения")) or "—",
        }
        photos = list_photos(item["id"], photos_dir)
        total_items += 1
        if not photos:
            cards_by_district.setdefault(district, []).append((item, None))
            continue
        for ph in photos:
            cards_by_district.setdefault(district, []).append((item, ph))
            total_photos += 1

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, cfg["brand"], len(cards_by_district), total_items, total_photos)

    # Сводная таблица — перед фотографиями
    build_summary_table_slide(prs, df, cfg)

    for district in sorted(cards_by_district.keys(), key=lambda s: s.lower()):
        cards = cards_by_district[district]
        unique_items = len({c[0]["id"] for c in cards})
        photo_count = sum(1 for c in cards if c[1] is not None)
        add_district_divider(prs, cfg["brand"], district, unique_items, photo_count)

        page_total = (len(cards) + MAX_PER_SLIDE - 1) // MAX_PER_SLIDE
        for page_idx in range(page_total):
            chunk = cards[page_idx * MAX_PER_SLIDE:(page_idx + 1) * MAX_PER_SLIDE]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_header(slide, cfg["brand"], district, photo_count)
            cells = cell_layout(len(chunk))
            for cell, (item, photo) in zip(cells, chunk):
                place_photo(slide, cell, photo, item)

    prs.save(output_pptx)

    # Проверка целостности: файл должен быть валидным zip и открываться как презентация
    if not zipfile.is_zipfile(output_pptx):
        raise RuntimeError("Сохранённый .pptx повреждён (не является zip-архивом)")
    try:
        Presentation(output_pptx)
    except Exception as e:
        raise RuntimeError(f"Сохранённый .pptx не открывается: {e}")

    print(f"Сохранено: {output_pptx}  (районов: {len(cards_by_district)}, "
          f"нарушений: {total_items}, фото: {total_photos})")
    return output_pptx


# ============================================================
#  ЧТЕНИЕ EXCEL + ТОЧКА ВХОДА
# ============================================================

def _read_excel_smart(excel_path: str) -> pd.DataFrame:
    with pd.ExcelFile(excel_path) as xls:
        sheet = SHEET_NAME_PREFERRED if SHEET_NAME_PREFERRED in xls.sheet_names else xls.sheet_names[0]
        print(f"Читаю лист «{sheet}» из {os.path.basename(excel_path)}")
        df = pd.read_excel(xls, sheet_name=sheet)
    return df


def process_tsafap_oati_file(excel_path: str, system: str = "ЦАФАП", *,
                             headless: bool = HEADLESS,
                             download_photos: bool = True) -> str:
    """
    Полный пайплайн ОАТИ/ЦАФАП. system = "ОАТИ" или "ЦАФАП".
    Возвращает путь к готовой презентации .pptx.
    Все промежуточные файлы — в одном временном каталоге (рядом с .pptx),
    который удаляется вызывающей стороной после отправки.
    """
    cfg = SYSTEM_CONFIG.get(system)
    if cfg is None:
        raise ValueError(f"Неизвестный режим: {system!r}. Ожидается 'ОАТИ' или 'ЦАФАП'.")

    run_dir = Path(tempfile.mkdtemp(prefix=f"{cfg['short'].lower()}_"))
    try:
        photos_dir = run_dir / "photos"
        photos_dir.mkdir(parents=True, exist_ok=True)

        ts = datetime.now().strftime("%d.%m.%Y_%H-%M")
        output_pptx = str(run_dir / f"Фото_{cfg['short']}_{ts}.pptx")

        df = _read_excel_smart(excel_path)
        print(f"Режим: {system}. Загружено строк: {len(df)}")

        if download_photos:
            download_phase(df, photos_dir, headless=headless)

        build_phase(df, photos_dir, output_pptx, cfg)
        return output_pptx
    except Exception:
        shutil.rmtree(run_dir, ignore_errors=True)
        raise


if __name__ == "__main__":
    import sys
    path = sys.argv[1] if len(sys.argv) > 1 else "1.xlsx"
    syst = sys.argv[2] if len(sys.argv) > 2 else "ЦАФАП"
    out = process_tsafap_oati_file(path, syst, headless=False)
    print("Готово:", out)